from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.security import generate_password_hash, check_password_hash
import jwt
from functools import wraps
import sqlite3
from datetime import datetime, timedelta
import os
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches
import google.generativeai as genai
from dotenv import load_dotenv
load_dotenv()
import io
import traceback
from flask import send_from_directory

# CRITICAL FIX: Since app.py is in /backend, point to ../frontend/build
BUILD_PATH = os.path.join(os.path.dirname(__file__), '..', 'frontend', 'build')

app = Flask(__name__, static_folder=BUILD_PATH, static_url_path='')

CORS(app, resources={
    r"/api/*": {
        "origins": ["https://infini-ai-doc-platform.vercel.app"]
    }
})

@app.route("/", defaults={"path": ""})
@app.route("/<path:path>")
def serve(path):
    if path != "" and os.path.exists(os.path.join(app.static_folder, path)):
        return send_from_directory(app.static_folder, path)
    else:
        return send_from_directory(app.static_folder, "index.html")

app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'your-secret-key-change-in-production')
app.config['DATABASE'] = os.environ.get('DATABASE_PATH', 'docgen.db')

#(User chose: gemini-2.0-flash-exp)
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY', '')
if not GEMINI_API_KEY:
    print("Warning: GEMINI_API_KEY not set. AI generation will fallback to placeholder content.")

try:
    genai.configure(api_key=GEMINI_API_KEY)
    MODEL_NAME = "gemini-2.0-flash"
    model = genai.GenerativeModel(model_name=MODEL_NAME)
except Exception as e:
    print("Warning: Failed to configure genai:", e)
    model = None
    MODEL_NAME = None


def get_db():
    db = sqlite3.connect(app.config['DATABASE'], timeout=10, check_same_thread=False)
    db.row_factory = sqlite3.Row
    return db

def init_db():
    db = get_db()
    try:
        db.executescript('''
            CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                email TEXT UNIQUE NOT NULL,
                password TEXT NOT NULL,
                name TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );

            CREATE TABLE IF NOT EXISTS projects (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER NOT NULL,
                document_type TEXT NOT NULL,
                title TEXT NOT NULL,
                topic TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (user_id) REFERENCES users (id)
            );

            CREATE TABLE IF NOT EXISTS sections (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                project_id INTEGER NOT NULL,
                title TEXT NOT NULL,
                content TEXT,
                order_index INTEGER NOT NULL,
                liked BOOLEAN,
                comment TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (project_id) REFERENCES projects (id)
            );

            CREATE TABLE IF NOT EXISTS refinement_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                section_id INTEGER NOT NULL,
                prompt TEXT NOT NULL,
                previous_content TEXT,
                new_content TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (section_id) REFERENCES sections (id)
            );
        ''')
        db.commit()
    finally:
        db.close()

#JWT
def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        token = request.headers.get('Authorization')
        
        if not token:
            return jsonify({'error': 'Token is missing'}), 401
        
        try:
            if token.startswith('Bearer '):
                token = token[7:]
            data = jwt.decode(token, app.config['SECRET_KEY'], algorithms=['HS256'])
            current_user_id = data.get('user_id')
            if not current_user_id:
                raise Exception("Invalid token payload")
        except Exception as e:
            return jsonify({'error': 'Token is invalid', 'details': str(e)}), 401
        
        return f(current_user_id, *args, **kwargs)
    
    return decorated


def _extract_text_from_response(response):
    """
    Robust extraction for many genai response shapes.
    Prefer: response.candidates[0].content.parts[0].text
    Fallback to other plausible fields.
    """
    try:
     
        if hasattr(response, "candidates") and response.candidates:
            c0 = response.candidates[0]
           
            if hasattr(c0, "content") and getattr(c0.content, "parts", None):
                part0 = c0.content.parts[0]
                if hasattr(part0, "text"):
                    return part0.text.strip()
                if isinstance(part0, dict) and "text" in part0:
                    return part0["text"].strip()

            if hasattr(c0, "text"):
                return c0.text.strip()
            if hasattr(c0, "content") and isinstance(c0.content, str):
                return c0.content.strip()

       
        if hasattr(response, "text"):
            return response.text.strip()

      
        return str(response).strip()
    except Exception:
        try:
            return str(response)
        except Exception:
            return ""

def _generate_with_model(prompt, max_output_tokens=512, temperature=0.2):
    """
    Correct Gemini (google-generativeai 0.3.2) call.
    """
    placeholder = "This is placeholder content. Configure your Gemini API key to generate AI content."

    if not GEMINI_API_KEY:
        return placeholder

    try:
        model = genai.GenerativeModel(MODEL_NAME)

        response = model.generate_content(
            prompt,
            generation_config={
                "max_output_tokens": max_output_tokens,
                "temperature": temperature
            }
        )

        # response.text is ALWAYS correct for Gemini 2.0
        if hasattr(response, "text") and response.text:
            return response.text.strip()

        return placeholder

    except Exception as e:
        print("AI CALL FAILED:", e)
        traceback.print_exc()
        return placeholder



def generate_content_with_ai(topic, section_title, document_type):
    try:
        if document_type == 'docx':
            prompt = f"""
Write detailed, high-quality content (200-300 words) for a document section.

Topic: {topic}
Section Title: {section_title}

Do NOT repeat the section title. Write only the content.
"""
            return _generate_with_model(prompt, max_output_tokens=600)

        else:  # pptx
            prompt = f"""
Create 4–6 bullet points for a PowerPoint slide.

Topic: {topic}
Slide Title: {section_title}

Format:
• Bullet point 1
• Bullet point 2
(1 sentence each)
"""
            return _generate_with_model(prompt, max_output_tokens=300)

    except:
        return "AI generation failed. Configure Gemini API properly."

def refine_content_with_ai(current_content, refinement_prompt, document_type):
    try:
        bullet_rule = "Keep bullet format with • symbols." if document_type == "pptx" else ""
        
        prompt = f"""
Refine the following content based on user request.

Original content:
{current_content}

User wants:
{refinement_prompt}

{bullet_rule}
Return ONLY the refined text.
"""

        return _generate_with_model(prompt, max_output_tokens=400)

    except:
        return current_content
def suggest_outline_with_ai(topic, document_type):
    try:
        if document_type == 'docx':
            prompt = f"""
Create 6–8 professional document section titles for topic:

{topic}

Return only titles, one per line, no numbering.
"""
        else:
            prompt = f"""
Create 8–12 PowerPoint slide titles for topic:

{topic}

Start with Title Slide, end with Thank You slide.
One title per line, no numbering.
"""

        text = _generate_with_model(prompt, max_output_tokens=200)
        return [line.strip() for line in text.split("\n") if line.strip()]

    except:
        if document_type == 'docx':
            return ["Introduction", "Background", "Analysis", "Findings", "Discussion", "Conclusion"]
        else:
            return ["Title Slide", "Introduction", "Overview", "Main Points", "Analysis", "Results", "Conclusion", "Thank You"]

@app.route('/api/auth/register', methods=['POST'])
def register():
    data = request.get_json() or {}
    
    email = data.get('email')
    password = data.get('password')
    name = data.get('name')
    
    if not email or not password or not name:
        return jsonify({'error': 'Missing required fields'}), 400
    
    db = get_db()
    try:
        # Check if user exists
        existing_user = db.execute('SELECT id FROM users WHERE email = ?', (email,)).fetchone()
        if existing_user:
            return jsonify({'error': 'Email already registered'}), 400
        
        hashed_password = generate_password_hash(password)
        cursor = db.execute(
            'INSERT INTO users (email, password, name) VALUES (?, ?, ?)',
            (email, hashed_password, name)
        )
        db.commit()
        user_id = cursor.lastrowid
        
        token = jwt.encode({
            'user_id': user_id,
            'exp': datetime.utcnow() + timedelta(days=30)
        }, app.config['SECRET_KEY'], algorithm='HS256')
        if isinstance(token, bytes):
            token = token.decode('utf-8')
        
        return jsonify({
            'token': token,
            'user': {
                'id': user_id,
                'email': email,
                'name': name
            }
        }), 201
    finally:
        db.close()

@app.route('/api/auth/login', methods=['POST'])
def login():
    data = request.get_json() or {}
    email = data.get('email')
    password = data.get('password')
    
    if not email or not password:
        return jsonify({'error': 'Missing email or password'}), 400
    
    db = get_db()
    try:
        user = db.execute('SELECT * FROM users WHERE email = ?', (email,)).fetchone()
        if not user or not check_password_hash(user['password'], password):
            return jsonify({'error': 'Invalid credentials'}), 401
        
        token = jwt.encode({
            'user_id': user['id'],
            'exp': datetime.utcnow() + timedelta(days=30)
        }, app.config['SECRET_KEY'], algorithm='HS256')
        if isinstance(token, bytes):
            token = token.decode('utf-8')
        
        return jsonify({
            'token': token,
            'user': {
                'id': user['id'],
                'email': user['email'],
                'name': user['name']
            }
        })
    finally:
        db.close()


@app.route('/api/projects', methods=['GET'])
@token_required
def get_projects(current_user_id):
    db = get_db()
    try:
        projects = db.execute(
            'SELECT * FROM projects WHERE user_id = ? ORDER BY updated_at DESC',
            (current_user_id,)
        ).fetchall()
        return jsonify({'projects': [dict(p) for p in projects]})
    finally:
        db.close()

@app.route('/api/projects', methods=['POST'])
@token_required
def create_project(current_user_id):
    data = request.get_json() or {}
    document_type = data.get('document_type')
    title = data.get('title')
    topic = data.get('topic')
    outline = data.get('outline', [])
    
    if not document_type or not title or not topic:
        return jsonify({'error': 'Missing required fields'}), 400
    
    if document_type not in ['docx', 'pptx']:
        return jsonify({'error': 'Invalid document type'}), 400
    
    db = get_db()
    try:
        cursor = db.execute(
            'INSERT INTO projects (user_id, document_type, title, topic) VALUES (?, ?, ?, ?)',
            (current_user_id, document_type, title, topic)
        )
        project_id = cursor.lastrowid
        
        for idx, section_title in enumerate(outline):
            db.execute(
                'INSERT INTO sections (project_id, title, order_index) VALUES (?, ?, ?)',
                (project_id, section_title, idx)
            )
        db.commit()
        
        project = db.execute('SELECT * FROM projects WHERE id = ?', (project_id,)).fetchone()
        return jsonify({'project': dict(project)}), 201
    finally:
        db.close()

@app.route('/api/projects/<int:project_id>', methods=['DELETE'])
@token_required
def delete_project(current_user_id, project_id):
    db = get_db()
    try:
        project = db.execute(
            'SELECT * FROM projects WHERE id = ? AND user_id = ?',
            (project_id, current_user_id)
        ).fetchone()
        if not project:
            return jsonify({'error': 'Project not found'}), 404
        
        db.execute('DELETE FROM refinement_history WHERE section_id IN (SELECT id FROM sections WHERE project_id = ?)', (project_id,))
        db.execute('DELETE FROM sections WHERE project_id = ?', (project_id,))
        db.execute('DELETE FROM projects WHERE id = ?', (project_id,))
        db.commit()
        return jsonify({'message': 'Project deleted'}), 200
    finally:
        db.close()

@app.route('/api/projects/<int:project_id>/sections', methods=['GET'])
@token_required
def get_sections(current_user_id, project_id):
    db = get_db()
    try:
        project = db.execute('SELECT * FROM projects WHERE id = ? AND user_id = ?', (project_id, current_user_id)).fetchone()
        if not project:
            return jsonify({'error': 'Project not found'}), 404
        
        sections = db.execute('SELECT * FROM sections WHERE project_id = ? ORDER BY order_index', (project_id,)).fetchall()
        return jsonify({'sections': [dict(s) for s in sections]})
    finally:
        db.close()

@app.route('/api/projects/<int:project_id>/generate', methods=['POST'])
@token_required
def generate_content(current_user_id, project_id):
    db = get_db()
    try:
        project = db.execute(
            'SELECT * FROM projects WHERE id = ? AND user_id = ?',
            (project_id, current_user_id)
        ).fetchone()

        if not project:
            return jsonify({'error': 'Project not found'}), 404

        sections = db.execute(
            'SELECT * FROM sections WHERE project_id = ? ORDER BY order_index',
            (project_id,)
        ).fetchall()

        updated_sections = []

        # AI generation outside DB lock
        for section in sections:
            content = generate_content_with_ai(
                topic=project['topic'],
                section_title=section['title'],
                document_type=project['document_type']
            )

            # Quick DB write (prevents locking)
            db.execute(
                'UPDATE sections SET content = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?',
                (content, section['id'])
            )

            updated = dict(section)
            updated['content'] = content
            updated_sections.append(updated)

        db.commit()

        return jsonify({'sections': updated_sections}), 200

    except Exception as e:
        print("GENERATE ERROR:", e)
        return jsonify({'error': 'Generation failed', 'details': str(e)}), 500

    finally:
        db.close()


@app.route('/api/sections/<int:section_id>/refine', methods=['POST'])
@token_required
def refine_section(current_user_id, section_id):
    data = request.get_json() or {}
    prompt = data.get('prompt')
    if not prompt:
        return jsonify({'error': 'Prompt is required'}), 400
    
    db = get_db()
    try:
        section = db.execute('''
            SELECT s.*, p.user_id, p.document_type 
            FROM sections s 
            JOIN projects p ON s.project_id = p.id 
            WHERE s.id = ?
        ''', (section_id,)).fetchone()
        if not section or section['user_id'] != current_user_id:
            return jsonify({'error': 'Section not found'}), 404
        
        previous_content = section['content'] or ""
        new_content = refine_content_with_ai(previous_content, prompt, section['document_type'])
        
        db.execute(
            'INSERT INTO refinement_history (section_id, prompt, previous_content, new_content) VALUES (?, ?, ?, ?)',
            (section_id, prompt, previous_content, new_content)
        )
        db.execute('UPDATE sections SET content = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?', (new_content, section_id))
        db.commit()
        return jsonify({'content': new_content})
    finally:
        db.close()

@app.route('/api/sections/<int:section_id>/feedback', methods=['POST'])
@token_required
def update_feedback(current_user_id, section_id):
    data = request.get_json() or {}
    liked = data.get('liked')
    
    db = get_db()
    try:
        section = db.execute('''
            SELECT s.* FROM sections s 
            JOIN projects p ON s.project_id = p.id 
            WHERE s.id = ? AND p.user_id = ?
        ''', (section_id, current_user_id)).fetchone()
        if not section:
            return jsonify({'error': 'Section not found'}), 404
        
        db.execute('UPDATE sections SET liked = ? WHERE id = ?', (liked, section_id))
        db.commit()
        return jsonify({'message': 'Feedback updated'})
    finally:
        db.close()

@app.route('/api/sections/<int:section_id>/comment', methods=['POST'])
@token_required
def update_comment(current_user_id, section_id):
    data = request.get_json() or {}
    comment = data.get('comment')
    
    db = get_db()
    try:
        section = db.execute('''
            SELECT s.* FROM sections s 
            JOIN projects p ON s.project_id = p.id 
            WHERE s.id = ? AND p.user_id = ?
        ''', (section_id, current_user_id)).fetchone()
        if not section:
            return jsonify({'error': 'Section not found'}), 404
        
        db.execute('UPDATE sections SET comment = ? WHERE id = ?', (comment, section_id))
        db.commit()
        return jsonify({'message': 'Comment updated'})
    finally:
        db.close()

@app.route('/api/ai/suggest-outline', methods=['POST'])
@token_required
def suggest_outline(current_user_id):
    data = request.get_json() or {}
    topic = data.get('topic')
    document_type = data.get('document_type')
    if not topic or not document_type:
        return jsonify({'error': 'Topic and document type required'}), 400
    
    outline = suggest_outline_with_ai(topic, document_type)
    return jsonify({'outline': outline})


# Export / Create DOCX / PPTX

@app.route('/api/projects/<int:project_id>/export', methods=['GET'])
@token_required
def export_document(current_user_id, project_id):
    db = get_db()
    try:
        project = db.execute('SELECT * FROM projects WHERE id = ? AND user_id = ?', (project_id, current_user_id)).fetchone()
        if not project:
            return jsonify({'error': 'Project not found'}), 404
        
        sections = db.execute('SELECT * FROM sections WHERE project_id = ? ORDER BY order_index', (project_id,)).fetchall()
        
        if project['document_type'] == 'docx':
            file_obj = create_docx(project, sections)
            mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        else:
            file_obj = create_pptx(project, sections)
            mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        # send_file accepts file-like object; ensure at 0
        file_obj.seek(0)
        return send_file(file_obj, mimetype=mimetype, as_attachment=True, download_name=f"{project['title']}.{project['document_type']}")
    finally:
        db.close()

def create_docx(project, sections):
    """Create a Word document and return BytesIO"""
    doc = Document()
    # Title (center)
    title = doc.add_heading(project['title'], 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    topic_para = doc.add_paragraph(project['topic'])
    topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if topic_para.runs:
        topic_para.runs[0].italic = True
    
    doc.add_paragraph()
    for section in sections:
        doc.add_heading(section['title'], level=1)
        if section['content']:
            # If content has bullet symbols, keep them as separate paragraphs
            content = section['content']
            lines = content.splitlines()
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                # preserve bullets or plain paragraphs
                doc.add_paragraph(line)
        else:
            doc.add_paragraph("No content generated yet.")
        doc.add_paragraph()
    
    out = io.BytesIO()
    doc.save(out)
    out.seek(0)
    return out

def create_pptx(project, sections):
    """Create a PPTX and return BytesIO"""
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    
    # Title slide using layout 0 (common)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = project['title']
    try:
        subtitle = slide.placeholders[1]
        subtitle.text = project['topic']
    except Exception:
        # some templates may not have placeholder[1]
        pass
    
    for section in sections:
        # Skip the very first if it's the title slide (order_index==0 often used)
        if section['order_index'] == 0:
            continue
        
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        title_shape = slide.shapes.title
        body_shape = None
        # find a placeholder for body
        try:
            body_shape = slide.placeholders[1]
        except Exception:
            # fallback: find a textbox
            for shp in slide.shapes:
                if shp.has_text_frame:
                    body_shape = shp
                    break
        
        if title_shape:
            title_shape.text = section['title']
        
        if body_shape and section['content']:
            text_frame = body_shape.text_frame
            text_frame.clear()
            lines = section['content'].splitlines()
            first_set = False
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                # Remove leading bullet symbols if exist
                while line and line[0] in '•-*–— ':
                    # strip leading bullet characters and whitespace
                    line = line.lstrip('•-*–— ').strip()
                if not first_set:
                    text_frame.text = line
                    first_set = True
                else:
                    p = text_frame.add_paragraph()
                    p.text = line
                    p.level = 0
        else:
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()
                tf.text = "No content generated yet."
    
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out



if __name__ == '__main__':
   
    init_db()
    app.run(debug=False, port=int(os.environ.get('PORT', 5000)))